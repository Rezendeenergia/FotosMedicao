import os
import io
import zipfile
import copy
import time
import logging
from concurrent.futures import ThreadPoolExecutor
from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
app = Flask(__name__, static_folder=BASE_DIR)
CORS(app, expose_headers=["X-Photos-Used", "X-Slides-Filled", "X-Barramentos", "X-Processing-Time"])

app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0

# Dimensões PADRÃO RETRATO: 7,62 x 10,16 cm
TARGET_W_CM = 7.62
TARGET_H_CM = 10.16
CM_TO_EMU = 360000
TARGET_W_EMU = int(TARGET_W_CM * CM_TO_EMU)
TARGET_H_EMU = int(TARGET_H_CM * CM_TO_EMU)

# Slots relatorio fotografico (4 fotos por slide)
PHOTO_SLOTS_4 = [
    {"name": "Retangulo 19", "x": 189290,  "y": 2478960, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},
    {"name": "Retangulo 41", "x": 3193696, "y": 2478960, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},
    {"name": "Retangulo 15", "x": 6198102, "y": 2478960, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},
    {"name": "Retangulo 8",  "x": 9202508, "y": 2478960, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},
]

# Slots base concretada — ordem: Poste | Barramento | Base
PHOTO_SLOTS_3 = [
    {"name": "Retângulo 41", "x": 4724400, "y": 2299850, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},  # Poste
    {"name": "Retângulo 19", "x": 720435,  "y": 2299851, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},  # Barramento
    {"name": "Retângulo 15", "x": 8728366, "y": 2299851, "cx": TARGET_W_EMU, "cy": TARGET_H_EMU},  # Base
]

BARRAMENTO_TEXTBOX = {
    "x": 5186217, "y": 1547498, "cx": 1819564, "cy": 369332
}

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

@app.route("/")
def index():
    return send_from_directory(BASE_DIR, "index.html")

@app.route("/<path:filename>")
def static_files(filename):
    return send_from_directory(BASE_DIR, filename)

@app.route("/health")
def health():
    return jsonify({"status": "ok", "message": "Backend Rezende Energia rodando!", "timestamp": time.time()})

def crop_and_resize(img_bytes, target_w_px, target_h_px):
    """
    Redimensiona mantendo toda a imagem visivel (sem corte).
    Preenche o espaco com a imagem esticada — sem bordas brancas.
    Texto do canto inferior direito sempre preservado.
    """
    img = Image.open(io.BytesIO(img_bytes))
    if img.mode != "RGB":
        img = img.convert("RGB")
    img_w, img_h = img.size

    # Pre-reduz para economizar RAM (max 1.5x o alvo)
    pre_w = int(target_w_px * 1.5)
    pre_h = int(target_h_px * 1.5)
    if img_w > pre_w or img_h > pre_h:
        img.thumbnail((pre_w, pre_h), Image.Resampling.BILINEAR)
        img_w, img_h = img.size

    # Estica a imagem para preencher exatamente o slot (sem corte, sem borda)
    img_resized = img.resize((target_w_px, target_h_px), Image.Resampling.BILINEAR)
    buf = io.BytesIO()
    img_resized.save(buf, format="JPEG", quality=65, subsampling=2)
    del img, img_resized  # libera RAM imediatamente
    return buf.getvalue()

def add_photo_to_slide(slide, slot, img_bytes, already_processed=False, is_landscape=None):
    """
    Insere foto no slot com orientacao inteligente:
    - Padrao: retrato 7,62 x 10,16 cm
    - Se a foto for paisagem (largura > altura): inverte slot para 10,16 x 7,62 cm
    Sem bordas brancas. Com efeito Predefinicao 1 do PowerPoint.
    """
    sp_tree = slide.shapes._spTree

    # Remove placeholder original — tolerante a acento (Retangulo vs Retângulo)
    import unicodedata
    def normaliza(s):
        return unicodedata.normalize('NFD', s).encode('ascii', 'ignore').decode('ascii').lower()
    slot_name_norm = normaliza(slot["name"])
    for sp in list(sp_tree):
        if sp.tag.split("}")[-1] == "sp":
            cNvPr = sp.find(".//" + qn("p:cNvPr"))
            if cNvPr is not None and normaliza(cNvPr.get("name", "")) == slot_name_norm:
                sp_tree.remove(sp)
                break

    # Verifica orientacao da foto
    if is_landscape is None:
        # Detecta da imagem original (nao foi pre-processada)
        try:
            probe = Image.open(io.BytesIO(img_bytes))
            is_landscape = probe.width > probe.height
        except Exception:
            is_landscape = False

    # Sempre usa as dimensoes exatas do slot (sem inverter, sem borda branca)
    # Foto esticada para preencher completamente o espaco definido
    final_cx = slot["cx"]
    final_cy = slot["cy"]
    final_x = slot["x"]
    final_y = slot["y"]

    if already_processed:
        img_resized = img_bytes  # ja foi pre-processado em paralelo
    else:
        target_w_px = int(final_cx / CM_TO_EMU * 2.54 * 96)
        target_h_px = int(final_cy / CM_TO_EMU * 2.54 * 96)
        img_resized = crop_and_resize(img_bytes, target_w_px, target_h_px)

    _, rId = slide.part.get_or_add_image_part(io.BytesIO(img_resized))

    # Efeito Predefinicao 1: sombra externa suave — sem linha de borda
    pic_xml = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="9999" name="{slot['name']}"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{rId}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="{final_x}" y="{final_y}"/>
      <a:ext cx="{final_cx}" cy="{final_cy}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln><a:noFill/></a:ln>
    <a:effectLst>
      <a:outerShdw blurRad="63500" dist="38100" dir="2700000" algn="tl" rotWithShape="0">
        <a:prstClr val="black">
          <a:alpha val="40000"/>
        </a:prstClr>
      </a:outerShdw>
    </a:effectLst>
  </p:spPr>
</p:pic>'''

    sp_tree.append(etree.fromstring(pic_xml))


def set_barramento_number(slide, numero):
    """
    Insere o numero do barramento no shape 'CaixaDeTexto 27'.
    Esse shape tem 0 runs — precisa criar o run via XML diretamente,
    copiando atributos do endParaRPr (negrito, lang, etc).
    """
    from lxml import etree

    NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for shape in slide.shapes:
        if shape.name == "CaixaDeTexto 27" and shape.has_text_frame:
            tf = shape.text_frame
            para = tf.paragraphs[0]
            p_elem = para._p

            # Pega atributos do endParaRPr para copiar na formatacao do run
            end_rpr = p_elem.find(f"{{{NSMAP_A}}}endParaRPr")

            # Cria elemento <a:r>
            r_elem = etree.SubElement(p_elem, f"{{{NSMAP_A}}}r")

            # Cria <a:rPr> copiando atributos de endParaRPr
            rpr = etree.SubElement(r_elem, f"{{{NSMAP_A}}}rPr")
            if end_rpr is not None:
                for attr, val in end_rpr.attrib.items():
                    rpr.set(attr, val)
            rpr.set("dirty", "0")

            # Move o <a:r> para antes do <a:endParaRPr>
            if end_rpr is not None:
                p_elem.remove(r_elem)
                end_rpr.addprevious(r_elem)

            # Cria <a:t> com o numero
            t_elem = etree.SubElement(r_elem, f"{{{NSMAP_A}}}t")
            t_elem.text = numero

            logger.info(f"Numero '{numero}' inserido em CaixaDeTexto 27")
            return

    logger.warning(f"CaixaDeTexto 27 nao encontrado no slide — numero '{numero}' nao inserido")


NS_R_EMU = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
RID_OFFSET = 100  # rIds do template sao deslocados +100 para liberar espaco para fotos novas

def remap_rids_in_element(element, id_map):
    """Renumera r:embed, r:link, r:id em todo o XML do slide via lxml."""
    for attr in [f'{{{NS_R_EMU}}}embed', f'{{{NS_R_EMU}}}link', f'{{{NS_R_EMU}}}id']:
        for el in element.iter():
            val = el.get(attr)
            if val and val in id_map:
                el.set(attr, id_map[val])


def duplicate_slide(prs, slide_index):
    """
    Duplica um slide de forma segura, sem conflito de rId.
    Renumera os rIds do template +100 para garantir que fotos novas
    nao colidam com os rIds ja existentes no XML duplicado.
    """
    from pptx.opc.packuri import PackURI
    from pptx.parts.slide import SlidePart
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    template = prs.slides[slide_index]
    template_part = template.part

    # Copia profunda do XML do slide
    new_element = copy.deepcopy(template_part._element)

    # Remapear rIds do template para +RID_OFFSET para liberar rId1..rId99 para fotos novas
    # Isso evita que get_or_add_image_part escolha um rId ja usado no XML do slide duplicado
    NOTES_RELTYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide'
    template_rels = template_part.rels
    old_to_new_rid = {}
    for old_rid in template_rels.keys():
        m = re.match(r'rId(\d+)', old_rid)
        if m:
            old_to_new_rid[old_rid] = f'rId{int(m.group(1)) + RID_OFFSET}'
    remap_rids_in_element(new_element, old_to_new_rid)

    # Partname unico
    existing_partnames = set(
        str(prs.slides[i].part.partname) for i in range(len(prs.slides))
    )
    idx = len(prs.slides) + 1
    while f'/ppt/slides/slide{idx}.xml' in existing_partnames:
        idx += 1
    new_partname = PackURI(f'/ppt/slides/slide{idx}.xml')

    # Cria o novo SlidePart
    new_part = SlidePart(
        new_partname,
        template_part.content_type,
        template_part.package,
        new_element
    )

    # Copia relacoes com os novos rIds remapeados
    for old_rid, rel in template_rels.items():
        if rel.reltype == NOTES_RELTYPE:
            continue  # nao copiar notas — causa corrupcao
        new_rid = old_to_new_rid.get(old_rid, old_rid)
        if rel.is_external:
            created = new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            created = new_part.relate_to(rel.target_part, rel.reltype, is_external=False)
        # Renomear rId criado pelo relate_to para o rId remapeado
        if created != new_rid:
            rel_obj = new_part._rels._rels.pop(created)
            new_part._rels._rels[new_rid] = rel_obj

    # Registra o novo slide na apresentacao
    slide_reltype = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
    rId = prs.slides.part.relate_to(new_part, slide_reltype)

    # Insere na lista XML de slides com id unico
    sldIdLst = prs.slides._sldIdLst
    max_id = max((int(s.get('id')) for s in sldIdLst), default=255)
    NS_PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldId = etree.SubElement(sldIdLst, f'{{{NS_PML}}}sldId')
    sldId.set('id', str(max_id + 1))
    sldId.set(f'{{{NS_REL}}}id', rId)

    return prs.slides[-1]


def remove_last_slide(prs):
    """Remove o último slide da apresentação."""
    sldIdLst = prs.slides._sldIdLst
    if len(sldIdLst) > 1:
        sldIdLst.remove(sldIdLst[-1])


def extract_photos_from_zip(zip_bytes):
    """Retorna lista de (nome, dados) — mantida para /validate e compatibilidade."""
    photos = []
    valid_ext = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}
    with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
        names = sorted([n for n in zf.namelist()
                        if not n.startswith('__MACOSX') and not os.path.basename(n).startswith('.')
                        and os.path.splitext(n.lower())[1] in valid_ext])
        for name in names:
            data = zf.read(name)
            if len(data) > 1000:
                photos.append((os.path.basename(name), data))
    return photos

def list_photo_names_in_zip(zip_bytes):
    """Retorna apenas os nomes das fotos validas, sem carregar dados na RAM."""
    valid_ext = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}
    with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
        names = sorted([n for n in zf.namelist()
                        if not n.startswith('__MACOSX') and not os.path.basename(n).startswith('.')
                        and os.path.splitext(n.lower())[1] in valid_ext
                        and zf.getinfo(n).file_size > 1000])
    return names


@app.route("/validate", methods=["POST"])
def validate():
    if 'zip' not in request.files:
        return jsonify({"error": "ZIP nao enviado"}), 400
    try:
        photos = extract_photos_from_zip(request.files['zip'].read())
        slides_needed = max(1, -(-len(photos) // 4))
        return jsonify({"photos": len(photos), "slides_needed": slides_needed, "names": [p[0] for p in photos[:20]]})
    except Exception as e:
        return jsonify({"error": str(e)}), 400


# Dimensoes alvo em pixels para as fotos (pre-calculado uma vez)
_TARGET_W_PX = int(TARGET_W_EMU / CM_TO_EMU * 2.54 * 96)
_TARGET_H_PX = int(TARGET_H_EMU / CM_TO_EMU * 2.54 * 96)

def _preprocess_photo(args):
    """
    Processa uma foto (crop+resize) sequencial.
    Sempre usa dimensoes fixas do slot W x H, sem inverter orientacao.
    Retorna (idx, img_bytes_processado, is_landscape=False).
    """
    idx, img_bytes = args
    try:
        data = crop_and_resize(img_bytes, _TARGET_W_PX, _TARGET_H_PX)
        return idx, data, False
    except Exception as e:
        logger.warning(f"Foto {idx} ignorada: {e}")
        return idx, None, False

@app.route("/process", methods=["POST"])
def process_pptx():
    t0 = time.time()
    if 'pptx' not in request.files or 'zip' not in request.files:
        return jsonify({"error": "Arquivos PPTX e ZIP sao obrigatorios"}), 400
    pptx_bytes = request.files['pptx'].read()
    zip_bytes = request.files['zip'].read()
    try:
        # 1) Lista nomes sem carregar dados — O(1) de RAM
        photo_names = list_photo_names_in_zip(zip_bytes)
        if not photo_names:
            return jsonify({"error": "Nenhuma foto encontrada no ZIP"}), 400

        n_photos = min(len(photo_names), 100)
        photo_names = photo_names[:n_photos]
        slides_needed = max(1, -(-n_photos // 4))
        logger.info(f"Processando {n_photos} fotos em {slides_needed} slides (a partir do slide 2)...")

        # 2) Monta apresentacao com slides corretos
        # Slide 0 = capa/template fixo, nao recebe fotos
        # Slides 1..N = slides de fotos (duplicados a partir do slide 1 do template)
        prs = Presentation(io.BytesIO(pptx_bytes))
        del pptx_bytes  # libera RAM do PPTX original
        if len(prs.slides) == 0:
            return jsonify({"error": "Apresentacao sem slides"}), 400

        # Slides 0 e 1 = fixos (capa + slide padrao), nao recebem fotos
        # Fotos comecam no slide 2 (indice 2)
        SLIDES_FIXOS = 2
        total_slides_needed = SLIDES_FIXOS + slides_needed

        # Duplica o slide 1 (indice 1) como template de fotos
        foto_template_idx = 1
        while len(prs.slides) < total_slides_needed:
            duplicate_slide(prs, foto_template_idx)
        while len(prs.slides) > total_slides_needed:
            remove_last_slide(prs)

        # 3) Processa e insere foto a foto direto do ZIP — nunca acumula tudo na RAM
        # Fotos vao para slides[2], slides[3], ... (slides 0 e 1 sao fixos)
        photos_used = 0
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            del zip_bytes  # libera bytes brutos do ZIP, mantendo apenas o objeto ZipFile
            for photo_idx, name in enumerate(photo_names):
                slide_idx = SLIDES_FIXOS + (photo_idx // 4)  # começa no slide 2
                slot_idx  = photo_idx % 4
                slide = prs.slides[slide_idx]
                slot  = PHOTO_SLOTS_4[slot_idx]
                try:
                    raw = zf.read(name)
                    _, data, lsc = _preprocess_photo((photo_idx, raw))
                    del raw  # libera original imediatamente
                    if data is not None:
                        add_photo_to_slide(slide, slot, data,
                                           already_processed=True, is_landscape=lsc)
                        del data
                        photos_used += 1
                except Exception as ex:
                    logger.warning(f"Foto {name} ignorada: {ex}")

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        elapsed = time.time() - t0
        logger.info(f"Concluido em {elapsed:.1f}s — {photos_used} fotos inseridas")
        response = send_file(out, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                             as_attachment=True, download_name='relatorio_preenchido.pptx')
        response.headers['X-Photos-Used'] = str(photos_used)
        response.headers['X-Slides-Filled'] = str(slides_needed)
        response.headers['X-Processing-Time'] = f"{elapsed:.2f}s"
        return response
    except Exception as e:
        logger.exception("Erro no processamento")
        return jsonify({"error": str(e)}), 500


@app.route("/process-base", methods=["POST"])
def process_base_concretada():
    """
    Base concretada: ordem Poste | Barramento | Base
    Dimensoes: 10,16 x 7,62 cm (paisagem, sem bordas brancas, com Predefinicao 1)
    """
    t0 = time.time()
    if 'pptx' not in request.files:
        return jsonify({"error": "Arquivo PPTX e obrigatorio"}), 400
    pptx_bytes = request.files['pptx'].read()
    numeros = request.form.getlist('numeros[]')
    if not numeros:
        return jsonify({"error": "Nenhum barramento enviado"}), 400
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        if len(prs.slides) == 0:
            return jsonify({"error": "Apresentacao sem slides"}), 400
        n_barramentos = len(numeros)
        # Slide 0 = capa fixa, nao recebe fotos
        # Barramentos comecam no slide 1 (segundo slide)
        if len(prs.slides) < 2:
            return jsonify({"error": "Template precisa ter pelo menos 2 slides"}), 400

        # Slide 0 = capa fixa, nunca tocada
        # Slide 1 = SEMPRE o template de barramento (limpo, sem fotos)
        SLIDES_FIXOS_BASE = 1
        total_slides_base = SLIDES_FIXOS_BASE + n_barramentos

        # Remove slides extras, mantendo so capa (0) + template (1)
        while len(prs.slides) > 2:
            remove_last_slide(prs)

        # Salva o template de barramento (slide 1) como PPTX de 1 slide separado
        # para garantir que cada duplicacao parte de um estado 100% limpo e isolado
        template_pptx = io.BytesIO()
        prs.save(template_pptx)
        template_pptx_bytes = template_pptx.getvalue()

        # Para cada barramento: carrega o template limpo e duplica slide 1 quantas vezes precisar
        # Isso evita que rels de imagens de um barramento vazem para o proximo
        prs_final = Presentation(io.BytesIO(template_pptx_bytes))

        # Garante que prs_final tem so capa + template (2 slides)
        while len(prs_final.slides) > 2:
            remove_last_slide(prs_final)

        # Duplica o template (slide 1) uma vez por barramento adicional
        for _ in range(n_barramentos - 1):
            duplicate_slide(prs_final, 1)

        # Agora insere fotos e numeros em cada slide — cada um eh uma copia limpa do template
        for i, numero in enumerate(numeros):
            slide = prs_final.slides[SLIDES_FIXOS_BASE + i]
            fotos = [
                request.files.get(f'poste_{i}'),
                request.files.get(f'barramento_{i}'),
                request.files.get(f'base_{i}'),
            ]
            for slot, foto in zip(PHOTO_SLOTS_3, fotos):
                if foto:
                    add_photo_to_slide(slide, slot, foto.read())
            if numero:
                set_barramento_number(slide, numero)
        prs = prs_final
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        elapsed = time.time() - t0
        response = send_file(out, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                             as_attachment=True, download_name='base_concretada_preenchida.pptx')
        response.headers['X-Barramentos'] = str(n_barramentos)
        response.headers['X-Processing-Time'] = f"{elapsed:.2f}s"
        return response
    except Exception as e:
        logger.exception("Erro no processamento de base")
        return jsonify({"error": str(e)}), 500


@app.route("/inspect-base", methods=["POST"])
def inspect_base():
    """
    Rota de diagnostico: recebe o pptx de base e retorna todos os shapes
    de todos os slides com nome, tipo e texto. Usar para identificar o
    campo do numero do barramento.
    """
    if 'pptx' not in request.files:
        return jsonify({"error": "Envie o arquivo pptx"}), 400
    pptx_bytes = request.files['pptx'].read()
    prs = Presentation(io.BytesIO(pptx_bytes))
    resultado = []
    for slide_i, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left_cm": round(shape.left / 360000, 2),
                "top_cm": round(shape.top / 360000, 2),
                "width_cm": round(shape.width / 360000, 2),
                "height_cm": round(shape.height / 360000, 2),
                "has_text": shape.has_text_frame,
                "text": shape.text_frame.text[:100] if shape.has_text_frame else None,
                "is_placeholder": shape.is_placeholder,
                "placeholder_idx": shape.placeholder_format.idx if shape.is_placeholder else None,
            }
            shapes_info.append(info)
        resultado.append({"slide": slide_i + 1, "shapes": shapes_info})
    return jsonify(resultado)


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
